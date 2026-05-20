#!/usr/bin/env python3
"""
PowerPoint Generator — Prediger Corporate Design
-------------------------------------------------
Erzeugt PPTX-Präsentationen aus einer JSON-Spezifikation.
Verwendet powerpoint_master.pptx als Vorlage; alle Folien nutzen
ausschließlich die vorhandenen Corporate-Design-Layouts.

Verwendung:
    python pptx_generator.py folien.json
    python pptx_generator.py folien.json --output meine_praesentation.pptx
    python pptx_generator.py folien.json --template /pfad/zur/powerpoint_master.pptx

JSON-Format: siehe Layout-Dokumentation in powerpoint_skill.json
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

DEFAULT_TEMPLATE = Path(__file__).parent / "powerpoint_master.pptx"

# Kürzel → Layout-Index in der Vorlage
LAYOUT_KEYS: dict[str, int] = {
    "titelfolie":       0,   # Präsentationstitel ohne Hintergrundbild
    "titelfolie_bild":  1,   # Präsentationstitel mit Hintergrundbild
    "kapitel_1":        2,   # Kapiteltrenner_1
    "kapitel_2":        3,   # Kapiteltrenner_2
    "agenda":           4,   # Agenda
    "inhalt":           5,   # Titel und Inhalt (großes Statement)
    "inhalt_bild":      6,   # Titel, Inhalt und Bild
    "zwei_spalten":     7,   # 2 Inhalte (2 Spalten)
    "drei_spalten":     8,   # 3 Inhalte (3 Spalten)
    "prozess":          9,   # Prozessablauf (5 Schritte mit Pfeilen)
    "zahl_gross":       11,  # Zahl groß (eine große Kennzahl)
    "vier_zahlen":      12,  # Vier Zahlen (4 Kennzahlen)
    "zwei_bilder":      15,  # 2 Bilder mit Untertitel
    "drei_bilder":      16,  # 3 Bilder mit Untertitel
    "bild_2_texte":     18,  # Bild mit 2 Texten
    "galerie_1":        19,  # Bildgalerie (6 Bilder)
    "galerie_2":        20,  # Bildgalerie (7 Bilder)
}


# ── Hilfsfunktionen ────────────────────────────────────────────────────────────

def remove_all_slides(prs: Presentation) -> None:
    """Entfernt alle Muster-Folien aus der geladenen Vorlage."""
    slide_id_list = prs.slides._sldIdLst
    for i in range(len(slide_id_list) - 1, -1, -1):
        r_id = slide_id_list[i].get(qn("r:id"))
        prs.part.drop_rel(r_id)
        slide_id_list.remove(slide_id_list[i])


def get_layout(prs: Presentation, key: str):
    """Gibt das Layout-Objekt für einen Schlüssel, Index oder Namen zurück."""
    if key in LAYOUT_KEYS:
        return prs.slide_layouts[LAYOUT_KEYS[key]]
    try:
        return prs.slide_layouts[int(key)]
    except (ValueError, IndexError):
        pass
    for layout in prs.slide_layouts:
        if layout.name == key:
            return layout
    valid = ", ".join(f'"{k}"' for k in LAYOUT_KEYS)
    raise ValueError(f"Unbekanntes Layout: '{key}'\nGültige Schlüssel: {valid}")


def resolve_layout_idx(prs: Presentation, key: str) -> int | None:
    """Gibt den numerischen Layout-Index zurück."""
    if key in LAYOUT_KEYS:
        return LAYOUT_KEYS[key]
    try:
        return int(key)
    except (ValueError, TypeError):
        pass
    for i, layout in enumerate(prs.slide_layouts):
        if layout.name == key:
            return i
    return None


def set_text(ph, text: str) -> None:
    """Schreibt einfachen Text in einen Platzhalter (löscht vorhandenen Inhalt)."""
    tf = ph.text_frame
    tf.clear()
    tf.paragraphs[0].add_run().text = str(text)


def set_bullets(ph, items: list) -> None:
    """Schreibt eine Liste als mehrzeiligen Inhalt in einen Platzhalter."""
    tf = ph.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.add_run().text = str(item)


def set_heading_and_bullets(ph, heading: str | None, bullets: list | None) -> None:
    """Schreibt Überschrift und optionale Stichpunkte in einen Platzhalter."""
    tf = ph.text_frame
    tf.clear()
    first = True
    if heading:
        tf.paragraphs[0].add_run().text = heading
        first = False
    if bullets:
        for item in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.add_run().text = str(item)
            first = False


def insert_image(ph, image_path: str) -> None:
    """Fügt ein Bild in einen Bildplatzhalter ein (optional)."""
    path = Path(image_path)
    if not path.exists():
        print(f"  ⚠️  Bilddatei nicht gefunden: {image_path}", file=sys.stderr)
        return
    try:
        ph.insert_picture(str(path))
    except Exception as e:
        print(f"  ⚠️  Bild konnte nicht eingefügt werden ({image_path}): {e}", file=sys.stderr)


# ── Folie befüllen ─────────────────────────────────────────────────────────────

def fill_slide(slide, spec: dict, layout_idx: int) -> None:
    """Füllt die Platzhalter einer Folie gemäß der Spezifikation."""
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}

    def t(idx: int, key: str, default=None):
        val = spec.get(key, default)
        if val is not None and idx in ph:
            set_text(ph[idx], val)

    def b(idx: int, key: str):
        val = spec.get(key)
        if val is not None and idx in ph:
            set_bullets(ph[idx], val)

    def img(idx: int, key: str):
        val = spec.get(key)
        if val and idx in ph:
            insert_image(ph[idx], val)

    if layout_idx == 0:
        # titelfolie: Präsentationstitel ohne Hintergrundbild
        # ph[0] = zentraler Titel
        t(0, "title")

    elif layout_idx == 1:
        # titelfolie_bild: Präsentationstitel mit Hintergrundbild
        # ph[0] = Titel, ph[10] = Hintergrundbild
        t(0, "title")
        img(10, "image")

    elif layout_idx in (2, 3):
        # kapitel_1 / kapitel_2: Kapiteltrenner
        # ph[11] = großer Kapiteltext
        t(11, "text")

    elif layout_idx == 4:
        # agenda: Agenda-Folie
        # ph[11] = Überschrift ("Agenda"), ph[12] = Themenliste
        t(11, "heading", "Agenda")
        b(12, "items")

    elif layout_idx == 5:
        # inhalt: Titel und Inhalt (großes Statement / Überschrift + optionale Stichpunkte)
        # ph[10] = kleine Oberüberschrift (Kategorie), ph[11] = Haupttext / Stichpunkte
        t(10, "supertitle")
        if 11 in ph:
            set_heading_and_bullets(ph[11], spec.get("heading"), spec.get("bullets"))

    elif layout_idx == 6:
        # inhalt_bild: Titel, Inhalt und Bild
        # ph[10] = Oberüberschrift, ph[12] = Text/Stichpunkte, ph[11] = Bild
        t(10, "supertitle")
        if 12 in ph:
            set_heading_and_bullets(ph[12], spec.get("heading"), spec.get("bullets"))
        img(11, "image")

    elif layout_idx == 7:
        # zwei_spalten: 2 Inhalte
        # ph[11] = Folientitel, ph[17]+ph[14] = Spalte 1, ph[19]+ph[18] = Spalte 2
        t(11, "heading")
        t(17, "col1_supertitle")
        b(14, "col1_bullets")
        t(19, "col2_supertitle")
        b(18, "col2_bullets")

    elif layout_idx == 8:
        # drei_spalten: 3 Inhalte
        # ph[11] = Folientitel, Spalten: (17,14) (19,18) (21,20)
        t(11, "heading")
        t(17, "col1_supertitle")
        b(14, "col1_bullets")
        t(19, "col2_supertitle")
        b(18, "col2_bullets")
        t(21, "col3_supertitle")
        b(20, "col3_bullets")

    elif layout_idx == 9:
        # prozess: Prozessablauf mit bis zu 5 Schritten
        # ph[11] = Überschrift
        # Schritt i: Label=ph[26+i], Titel=ph[17,19,21,23,25], Inhalt=ph[14,18,20,22,24]
        t(11, "heading")
        steps = spec.get("steps", [])
        label_idx   = [26, 27, 28, 29, 30]
        title_idx   = [17, 19, 21, 23, 25]
        content_idx = [14, 18, 20, 22, 24]
        for i, step in enumerate(steps[:5]):
            if label_idx[i] in ph:
                set_text(ph[label_idx[i]], step.get("label", f"Schritt {i + 1}"))
            if title_idx[i] in ph:
                set_text(ph[title_idx[i]], step.get("title", ""))
            if content_idx[i] in ph:
                content = step.get("bullets") or step.get("text")
                if isinstance(content, list):
                    set_bullets(ph[content_idx[i]], content)
                elif content:
                    set_text(ph[content_idx[i]], content)

    elif layout_idx == 11:
        # zahl_gross: eine große Kennzahl
        # ph[11] = Zahl/Wert, ph[12] = Beschreibung
        t(11, "number")
        t(12, "description")

    elif layout_idx == 12:
        # vier_zahlen: vier Kennzahlen nebeneinander
        # ph[11,13,15,17] = Werte, ph[12,14,16,18] = Beschreibungen
        numbers = spec.get("numbers", [])
        val_idx  = [11, 13, 15, 17]
        desc_idx = [12, 14, 16, 18]
        for i, entry in enumerate(numbers[:4]):
            if val_idx[i] in ph:
                set_text(ph[val_idx[i]], entry.get("value", ""))
            if desc_idx[i] in ph:
                set_text(ph[desc_idx[i]], entry.get("description", ""))

    elif layout_idx == 15:
        # zwei_bilder: 2 Bilder mit Untertitel
        # ph[11] = Folientitel, ph[1]+ph[12] = Bild1+Caption, ph[15]+ph[16] = Bild2+Caption
        t(11, "heading")
        images = spec.get("images", [])
        img_idx     = [1, 15]
        caption_idx = [12, 16]
        for i, entry in enumerate(images[:2]):
            if isinstance(entry, dict):
                if img_idx[i] in ph:
                    insert_image(ph[img_idx[i]], entry.get("image", ""))
                if caption_idx[i] in ph:
                    set_text(ph[caption_idx[i]], entry.get("caption", ""))
            elif isinstance(entry, str) and img_idx[i] in ph:
                insert_image(ph[img_idx[i]], entry)

    elif layout_idx == 16:
        # drei_bilder: 3 Bilder mit Untertitel
        # ph[11] = Folientitel, (ph[1],ph[12]) (ph[13],ph[14]) (ph[15],ph[16])
        t(11, "heading")
        images = spec.get("images", [])
        img_idx     = [1, 13, 15]
        caption_idx = [12, 14, 16]
        for i, entry in enumerate(images[:3]):
            if isinstance(entry, dict):
                if img_idx[i] in ph:
                    insert_image(ph[img_idx[i]], entry.get("image", ""))
                if caption_idx[i] in ph:
                    set_text(ph[caption_idx[i]], entry.get("caption", ""))
            elif isinstance(entry, str) and img_idx[i] in ph:
                insert_image(ph[img_idx[i]], entry)

    elif layout_idx == 18:
        # bild_2_texte: Bild mit 2 Textspalten
        # ph[11] = Folientitel, ph[19] = Bild, ph[12]+ph[18] = Text links, ph[20]+ph[21] = Text rechts
        t(11, "heading")
        img(19, "image")
        t(12, "left_supertitle")
        t(18, "left_text")
        t(20, "right_text")
        t(21, "right_supertitle")


# ── Präsentation aufbauen ──────────────────────────────────────────────────────

def build_presentation(spec: dict, template_path: Path) -> Presentation:
    prs = Presentation(str(template_path))
    remove_all_slides(prs)

    slides = spec.get("slides", [])
    for i, slide_spec in enumerate(slides, 1):
        layout_key = slide_spec.get("layout", "inhalt")
        try:
            layout = get_layout(prs, layout_key)
        except ValueError as e:
            print(f"  ⚠️  Folie {i}: {e}", file=sys.stderr)
            continue

        layout_idx = resolve_layout_idx(prs, layout_key)
        slide = prs.slides.add_slide(layout)
        fill_slide(slide, slide_spec, layout_idx)
        print(f"  ✓ Folie {i:>2}: [{layout_key}] {layout.name}")

    return prs


def main():
    parser = argparse.ArgumentParser(description="PPTX aus JSON-Spezifikation erzeugen")
    parser.add_argument("input", help="JSON-Spezifikationsdatei")
    parser.add_argument("--output", "-o", help="Ausgabedatei (.pptx)")
    parser.add_argument(
        "--template", "-t",
        default=str(DEFAULT_TEMPLATE),
        help=f"Vorlage (Standard: {DEFAULT_TEMPLATE.name})",
    )
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"❌ Eingabedatei nicht gefunden: {input_path}", file=sys.stderr)
        sys.exit(1)

    with open(input_path, encoding="utf-8") as f:
        spec = json.load(f)

    output_path = Path(args.output) if args.output else input_path.with_suffix(".pptx")
    template_path = Path(args.template)

    if not template_path.exists():
        print(f"❌ Vorlage nicht gefunden: {template_path}", file=sys.stderr)
        sys.exit(1)

    slide_count = len(spec.get("slides", []))
    print(f"\n📊 Erzeuge: '{spec.get('title', output_path.stem)}'")
    print(f"   Vorlage:  {template_path}")
    print(f"   Ausgabe:  {output_path}")
    print(f"   Folien:   {slide_count}\n")

    prs = build_presentation(spec, template_path)
    prs.save(str(output_path))

    print(f"\n✅ Gespeichert: {output_path} ({slide_count} Folien)")


if __name__ == "__main__":
    main()
